"""Generate a demo .pptx so the pipeline can be exercised without real assets.

Deliberately uses theme colours, shape fills, a table and rotation — the slide
renderers are only meaningfully tested by a deck that has styling to lose.

    uv run python scripts/make_demo.py tmp/demo.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

INK = RGBColor(0x1C, 0x24, 0x33)
ACCENT = RGBColor(0xC4, 0x6A, 0x46)
MUTED = RGBColor(0x5B, 0x64, 0x72)
PANEL = RGBColor(0xF4, 0xF1, 0xEA)

SLIDES = [
    ("Doc2Video Agent", "把 PDF / PPT 变成会讲、会指、会切镜头的讲解视频", []),
    ("目录", "", ["产品定位", "系统架构", "核心指标", "落地案例", "下一步"]),
    (
        "产品定位",
        "面向 PDF/PPT 的 AI 视频导演插件",
        ["上传文档 + 一句话需求", "自动讲稿、配音、镜头", "对话式局部修改"],
    ),
    (
        "系统架构",
        "Agent 编排 Skills 与 Tools",
        [
            "Document Skill 解析文档",
            "Narration Skill 生成讲稿",
            "Director Skill 设计镜头",
            "确定性渲染器出片",
        ],
    ),
    ("核心指标", "上线三个月的关键数据", ["月活 12000", "平均生成时长 6.5 分钟", "复用率 78%"]),
    ("增长趋势", "四个季度的用量与留存", []),
    ("用量构成", "各类文档的占比", []),
    (
        "总结",
        "先证明一句话就能得到一条专业讲解视频",
        ["MVP 不做数字人", "信息型画面坚持确定性渲染"],
    ),
]


def build(out_path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for slide_index, (title, subtitle, bullets) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)

        # An accent band across the top: gives every slide a fill to render.
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = ACCENT
        band.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(1.1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_run = title_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(40)
        title_run.font.bold = True
        title_run.font.color.rgb = INK

        top = 2.0
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.4), Inches(0.7))
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_run = sub_frame.paragraphs[0].runs[0]
            sub_run.font.size = Pt(22)
            sub_run.font.color.rgb = MUTED
            top += 0.9

        if bullets:
            panel = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8),
                Inches(top),
                Inches(8.4),
                Inches(0.62 * len(bullets) + 0.5),
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = PANEL
            panel.line.color.rgb = ACCENT
            panel.line.width = Pt(1.5)
            # Shapes carry their own text frame; keep it empty and overlay text,
            # so bullet geometry stays independent of the panel.
            panel.text_frame.text = ""

            body = slide.shapes.add_textbox(
                Inches(1.1), Inches(top + 0.22), Inches(7.8), Inches(0.62 * len(bullets))
            )
            frame = body.text_frame
            frame.word_wrap = True
            for index, bullet in enumerate(bullets):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 1
                run = paragraph.runs[0]
                run.font.size = Pt(20)
                run.font.color.rgb = INK

        # Page number, right-aligned — exercises alignment and small text.
        footer = slide.shapes.add_textbox(Inches(8.2), Inches(6.7), Inches(1.2), Inches(0.4))
        footer_frame = footer.text_frame
        footer_frame.text = f"{slide_index} / {len(SLIDES)}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.RIGHT
        footer_para.runs[0].font.size = Pt(12)
        footer_para.runs[0].font.color.rgb = MUTED

        if title == "核心指标":
            _add_table(slide, top + 0.62 * len(bullets) + 0.5)
        elif title == "增长趋势":
            _add_column_chart(slide, top)
        elif title == "用量构成":
            _add_pie_chart(slide, top)

        notes = slide.notes_slide.notes_text_frame
        notes.text = f"讲这一页时强调：{subtitle or title}"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"已生成 {out_path}（{len(SLIDES)} 页）")


def _add_table(slide, top_inches: float) -> None:
    rows, cols = 3, 3
    shape = slide.shapes.add_table(
        rows, cols, Inches(0.8), Inches(top_inches), Inches(8.4), Inches(1.2)
    )
    table = shape.table
    data = [
        ["指标", "本月", "环比"],
        ["月活", "12000", "+18%"],
        ["复用率", "78%", "+6pt"],
    ]
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.runs[0].font.size = Pt(14)
            if r == 0:
                paragraph.runs[0].font.bold = True


def _add_column_chart(slide, top_inches: float) -> None:
    """Clustered column with two series — exercises grouping, legend, gridlines."""
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3", "Q4"]
    data.add_series("生成视频数", (1200, 2100, 3400, 5200))
    data.add_series("复用次数", (300, 780, 1600, 3100))

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(top_inches),
        Inches(8.4),
        Inches(3.8),
        data,
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


def _add_pie_chart(slide, top_inches: float) -> None:
    data = CategoryChartData()
    data.categories = ["PPTX", "PDF", "PPT"]
    data.add_series("占比", (0.62, 0.31, 0.07))

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2.2),
        Inches(top_inches),
        Inches(5.6),
        Inches(3.8),
        data,
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("tmp/demo.pptx")
    build(target)
