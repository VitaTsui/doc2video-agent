"""Placeholder inheritance and table styles.

Both read parts of the package python-pptx does not surface, and both decide
what a slide actually *looks* like rather than what it says.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from doc2video.tools.slides.extract import extract_deck
from doc2video.tools.slides.inherit import StyleResolver
from doc2video.tools.slides.table_style import load_table_styles
from doc2video.tools.slides.theme import Theme

# --------------------------------------------------------------------------
# placeholder inheritance
# --------------------------------------------------------------------------


@pytest.fixture
def layout_deck(tmp_path: Path) -> Path:
    """A deck built the ordinary way: standard layout, no formatting on the slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "标题"
    body = slide.placeholders[1].text_frame
    body.text = "一级"
    for text, level in (("二级", 1), ("三级", 2)):
        paragraph = body.add_paragraph()
        paragraph.text = text
        paragraph.level = level
    path = tmp_path / "layout.pptx"
    prs.save(path)
    return path


def test_resolver_reads_master_sizes_per_level(layout_deck: Path):
    prs = Presentation(layout_deck)
    resolver = StyleResolver(prs, Theme())
    body = prs.slides[0].placeholders[1]

    sizes = [resolver.defaults_for(body, level).size_pt for level in range(3)]
    assert all(size is not None for size in sizes), f"每一层都应继承到字号：{sizes}"
    assert sizes[0] > sizes[1] > sizes[2]


def test_resolver_reads_master_bullets(layout_deck: Path):
    prs = Presentation(layout_deck)
    resolver = StyleResolver(prs, Theme())
    body = prs.slides[0].placeholders[1]

    assert resolver.defaults_for(body, 0).bullet, "正文占位符应从母版继承项目符号"


def test_resolver_reads_title_alignment(layout_deck: Path):
    prs = Presentation(layout_deck)
    resolver = StyleResolver(prs, Theme())
    title = prs.slides[0].shapes.title

    defaults = resolver.defaults_for(title, 0)
    assert defaults.size_pt == pytest.approx(44.0)
    assert defaults.align == "center"


def test_plain_text_box_inherits_no_bullet(tmp_path: Path):
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    box.text_frame.text = "纯文本框"
    path = tmp_path / "plain.pptx"
    prs.save(path)

    prs2 = Presentation(path)
    resolver = StyleResolver(prs2, Theme())
    shape = next(s for s in prs2.slides[0].shapes if s.has_text_frame)
    assert not resolver.defaults_for(shape, 0).bullet


def test_extraction_applies_inheritance(layout_deck: Path, tmp_path: Path):
    deck = extract_deck(layout_deck, tmp_path / "assets", target_width=1920)
    bodies = [
        shape
        for shape in deck.slides[0].shapes
        if shape.text and len(shape.text.paragraphs) == 3
    ]
    assert bodies, "应当找到三段正文的占位符"

    sizes = [para.runs[0].size_pt for para in bodies[0].text.paragraphs]
    assert all(size is not None for size in sizes)
    assert sizes[0] > sizes[1] > sizes[2]
    assert all(para.bullet for para in bodies[0].text.paragraphs)


def test_font_scale_shrinks_inherited_sizes(layout_deck: Path, tmp_path: Path):
    """PowerPoint's autofit stores a shrink factor that must reach the render."""
    from pptx.oxml.ns import qn

    prs = Presentation(layout_deck)
    body = prs.slides[0].placeholders[1]
    body_pr = body.text_frame._txBody.find(qn("a:bodyPr"))
    norm = body_pr.makeelement(qn("a:normAutofit"), {"fontScale": "50000"})
    body_pr.append(norm)
    scaled = tmp_path / "scaled.pptx"
    prs.save(scaled)

    plain = extract_deck(layout_deck, tmp_path / "a1", target_width=1920)
    shrunk = extract_deck(scaled, tmp_path / "a2", target_width=1920)

    def first_body_size(deck):
        for shape in deck.slides[0].shapes:
            if shape.text and len(shape.text.paragraphs) == 3:
                return shape.text.paragraphs[0].runs[0].size_pt
        raise AssertionError("找不到正文占位符")

    assert first_body_size(shrunk) == pytest.approx(first_body_size(plain) * 0.5)


# --------------------------------------------------------------------------
# table styles
# --------------------------------------------------------------------------

TABLE_STYLES_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               def="{11111111-1111-1111-1111-111111111111}">
  <a:tblStyle styleId="{11111111-1111-1111-1111-111111111111}" styleName="Demo">
    <a:wholeTbl>
      <a:tcTxStyle><a:srgbClr val="203040"/></a:tcTxStyle>
      <a:tcStyle>
        <a:tcBdr>
          <a:left><a:ln><a:solidFill><a:srgbClr val="AABBCC"/></a:solidFill></a:ln></a:left>
        </a:tcBdr>
        <a:fill><a:solidFill><a:srgbClr val="F7F7F7"/></a:solidFill></a:fill>
      </a:tcStyle>
    </a:wholeTbl>
    <a:band1H>
      <a:tcStyle><a:fill><a:solidFill><a:srgbClr val="E4ECF7"/></a:solidFill></a:fill></a:tcStyle>
    </a:band1H>
    <a:firstRow>
      <a:tcTxStyle><a:srgbClr val="FFFFFF"/></a:tcTxStyle>
      <a:tcStyle>
        <a:fill><a:solidFill><a:schemeClr val="accent2"/></a:solidFill></a:fill>
      </a:tcStyle>
    </a:firstRow>
  </a:tblStyle>
</a:tblStyleLst>
"""


@pytest.fixture
def styled_package(tmp_path: Path) -> Path:
    path = tmp_path / "styles.pptx"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("ppt/tableStyles.xml", TABLE_STYLES_XML)
    return path


def test_table_style_is_read_from_the_package(styled_package: Path):
    styles = load_table_styles(styled_package, Theme())
    style = styles.get("{11111111-1111-1111-1111-111111111111}")

    assert style is not None
    assert style.header_fill == "#ED7D31"  # accent2 resolved through the theme
    assert style.header_text == "#FFFFFF"
    assert style.band_fill == "#E4ECF7"
    assert style.body_fill == "#F7F7F7"
    assert style.border_color == "#AABBCC"


def test_unknown_style_id_falls_back_to_the_declared_default(styled_package: Path):
    styles = load_table_styles(styled_package, Theme())
    assert styles.get("{does-not-exist}") is not None


def test_missing_part_yields_no_styles(tmp_path: Path):
    empty = tmp_path / "empty.pptx"
    with zipfile.ZipFile(empty, "w") as archive:
        archive.writestr("ppt/presentation.xml", "<x/>")

    styles = load_table_styles(empty, Theme())
    assert styles.get(None) is None


def test_extraction_keeps_the_accent_approximation_without_definitions(
    demo_pptx: Path, tmp_path: Path
):
    """The demo deck declares a style id but ships no definition for it."""
    deck = extract_deck(demo_pptx, tmp_path / "assets", target_width=1920)
    tables = [s.table for slide in deck.slides for s in slide.shapes if s.table]
    assert tables
    # Falls back rather than rendering a bare grid.
    assert tables[0].header_fill
    assert tables[0].band_fill


def test_the_style_cache_pins_the_objects_it_keys_on(demo_pptx: Path):
    """The cache is keyed by ``id()``, which is only unique while the object lives.

    python-pptx builds its element proxies on demand, so an entry whose key
    object has been collected can be hit by a *different* element allocated at
    the same address — handing that shape the previous shape's inherited size,
    colour and bullet. Each entry therefore keeps a reference to its key object.
    Dropping that reference reintroduces a bug that only shows up as an
    occasional wrong bullet in a rendered slide.
    """
    from pptx import Presentation

    from doc2video.tools.slides.inherit import StyleResolver
    from doc2video.tools.slides.theme import load_theme

    presentation = Presentation(str(demo_pptx))
    resolver = StyleResolver(presentation, load_theme(demo_pptx))
    for slide in presentation.slides:
        for shape in slide.shapes:
            resolver.defaults_for(shape, 0)

    assert resolver._shape_cache, "示例 deck 应至少缓存一个形状"
    for key, entry in resolver._shape_cache.items():
        pinned, _value = entry
        assert id(pinned) == key, "缓存必须持有键所对应的对象，否则地址会被复用"
