"""Charts python-pptx reports wrongly, or refuses to open at all.

``chart.chart_type`` names one type for the whole chart, which is wrong for a
combo, and ``PlotFactory`` raises outright for the 3D tags — so both used to
come out of the extractor as something other than what the slide shows.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from doc2video.tools.slides.chart_xml import read_plot_groups
from doc2video.tools.slides.extract import extract_deck
from doc2video.tools.slides.model import ChartKind
from doc2video.tools.slides.theme import Theme

C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

SECONDARY_MAX = 0.5
REVENUE = (1200.0, 2400.0, 3600.0, 5200.0)
MARGIN = (0.31, 0.34, 0.38, 0.41)


def _chart_of(deck, index: int = 0):
    for shape in deck.slides[index].shapes:
        if shape.chart:
            return shape.chart
    raise AssertionError("这一页没有提取出图表")


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


@pytest.fixture
def combo_deck(tmp_path: Path) -> Path:
    """Columns plus a line on a second value axis — a margin-over-revenue chart."""
    prs = Presentation()
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3", "Q4"]
    data.add_series("营收", REVENUE)
    data.add_series("毛利率", MARGIN)
    chart = _blank(prs).shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(0.5), Inches(6), Inches(4),
        data,
    ).chart

    plot_area = chart._chartSpace.find(f"{C}chart/{C}plotArea")
    series = plot_area.find(f"{C}barChart").findall(f"{C}ser")

    line = etree.SubElement(plot_area, f"{C}lineChart")
    etree.SubElement(line, f"{C}grouping").set("val", "standard")
    line.append(series[1])
    marker = etree.SubElement(series[1], f"{C}marker")
    etree.SubElement(marker, f"{C}symbol").set("val", "circle")
    for axis_id in ("111111111", "222222222"):
        etree.SubElement(line, f"{C}axId").set("val", axis_id)

    value_axis = etree.SubElement(plot_area, f"{C}valAx")
    etree.SubElement(value_axis, f"{C}axId").set("val", "222222222")
    scaling = etree.SubElement(value_axis, f"{C}scaling")
    etree.SubElement(scaling, f"{C}max").set("val", str(SECONDARY_MAX))
    etree.SubElement(scaling, f"{C}min").set("val", "0")
    etree.SubElement(value_axis, f"{C}delete").set("val", "0")

    path = tmp_path / "combo.pptx"
    prs.save(path)
    return path


@pytest.fixture
def pie_3d_deck(tmp_path: Path) -> Path:
    """A 3D pie: the tag python-pptx will not open."""
    prs = Presentation()
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("份额", (40, 35, 25))
    chart = _blank(prs).shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(0.5), Inches(0.5), Inches(5), Inches(4), data
    ).chart
    plot_area = chart._chartSpace.find(f"{C}chart/{C}plotArea")
    plot_area.find(f"{C}pieChart").tag = f"{C}pie3DChart"

    path = tmp_path / "pie3d.pptx"
    prs.save(path)
    return path


# --------------------------------------------------------------------------
# combo charts and the secondary axis
# --------------------------------------------------------------------------


def test_line_series_is_not_drawn_as_a_bar(combo_deck: Path, tmp_path: Path):
    chart = _chart_of(extract_deck(combo_deck, tmp_path / "a", target_width=1920))

    assert chart.kind is ChartKind.COLUMN
    assert chart.series[0].kind is None, "第一条系列与图表类型相同，不必重复标注"
    assert chart.series[1].kind is ChartKind.LINE_MARKERS


def test_second_series_is_marked_as_secondary(combo_deck: Path, tmp_path: Path):
    chart = _chart_of(extract_deck(combo_deck, tmp_path / "a", target_width=1920))

    assert not chart.series[0].secondary_axis
    assert chart.series[1].secondary_axis
    assert chart.y2_max == pytest.approx(SECONDARY_MAX)
    assert chart.y2_visible


def test_primary_bounds_do_not_come_from_the_secondary_axis(
    combo_deck: Path, tmp_path: Path
):
    """python-pptx answers `value_axis` with *an* axis, and picks the wrong one.

    Taking its bounds pinned every bar to the top of a chart scaled 0–0.5.
    """
    chart = _chart_of(extract_deck(combo_deck, tmp_path / "a", target_width=1920))

    assert chart.y_max is None or chart.y_max > max(REVENUE) / 2


def test_a_plain_chart_gains_no_secondary_axis(demo_pptx: Path, tmp_path: Path):
    deck = extract_deck(demo_pptx, tmp_path / "a", target_width=1920)
    charts = [s.chart for slide in deck.slides for s in slide.shapes if s.chart]
    assert charts

    for chart in charts:
        assert not any(s.secondary_axis for s in chart.series)
        assert chart.y2_min is None and chart.y2_max is None
        assert not chart.three_d


# --------------------------------------------------------------------------
# plot types python-pptx cannot open
# --------------------------------------------------------------------------


def test_three_d_pie_is_extracted_rather_than_dropped(pie_3d_deck: Path, tmp_path: Path):
    """PlotFactory raises on `pie3DChart`, which used to lose the whole chart."""
    chart = _chart_of(extract_deck(pie_3d_deck, tmp_path / "a", target_width=1920))

    assert chart.kind is ChartKind.PIE
    assert chart.three_d
    assert chart.categories == ["A", "B", "C"]
    assert chart.series[0].values == [40, 35, 25]


def test_series_name_survives_the_xml_fallback(pie_3d_deck: Path, tmp_path: Path):
    chart = _chart_of(extract_deck(pie_3d_deck, tmp_path / "a", target_width=1920))
    assert chart.series[0].name == "份额"


def test_gaps_in_the_data_stay_gaps(tmp_path: Path):
    """A missing point must not become a zero — the renderer draws through it."""
    prs = Presentation()
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("s", (10, None, 30))
    chart = _blank(prs).shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(0.5), Inches(0.5), Inches(5), Inches(4), data
    ).chart
    plot_area = chart._chartSpace.find(f"{C}chart/{C}plotArea")
    plot_area.find(f"{C}lineChart").tag = f"{C}line3DChart"
    path = tmp_path / "gap.pptx"
    prs.save(path)

    groups = read_plot_groups(Presentation(path).slides[0].shapes[0].chart, Theme())
    assert groups.series[0].values == [10.0, None, 30.0]


def test_stacked_grouping_is_read_from_the_plot_group(tmp_path: Path):
    prs = Presentation()
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("s1", (1, 2))
    data.add_series("s2", (3, 4))
    chart = _blank(prs).shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(0.5), Inches(5), Inches(4), data
    ).chart
    groups = read_plot_groups(chart, Theme())

    assert groups.kinds == [ChartKind.COLUMN_STACKED, ChartKind.COLUMN_STACKED]
    assert not groups.has_secondary
