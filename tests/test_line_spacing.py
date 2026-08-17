"""Line spacing: PowerPoint's two kinds, both as a CSS line-height multiple.

The bug this guards against made whole paragraphs disappear from rendered
slides while looking perfectly correct everywhere else: the text was extracted,
positioned and styled, and then pushed millions of pixels below its own box by
a line-height of 266700. Absolute spacing (`<a:spcPts>`) arrives as a
``Length``, which subclasses ``int`` — so the branch written to catch it never
fired and raw EMU reached the renderer as a multiple.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

from doc2video.tools.slides.extract import extract_deck


def _deck(tmp_path: Path, *, spacing_xml: str | None) -> Path:
    """A slide whose body paragraph carries the given `<a:lnSpc>`, if any."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "标题"
    body = slide.placeholders[1].text_frame
    body.text = "正文一行"
    body.paragraphs[0].runs[0].font.size = Pt(13.5)

    if spacing_xml is not None:
        p_pr = body.paragraphs[0]._p.get_or_add_pPr()
        lnSpc = p_pr.makeelement(qn("a:lnSpc"), {})
        child_tag, val = spacing_xml
        child = lnSpc.makeelement(qn(child_tag), {"val": val})
        lnSpc.append(child)
        p_pr.insert(0, lnSpc)

    path = tmp_path / "spacing.pptx"
    prs.save(path)
    return path


def _body_spacing(deck_path: Path, assets: Path) -> float:
    deck = extract_deck(deck_path, assets, target_width=1920)
    for shape in deck.slides[0].shapes:
        if not shape.text:
            continue
        for paragraph in shape.text.paragraphs:
            if any("正文" in run.text for run in paragraph.runs):
                return paragraph.line_spacing
    raise AssertionError("没找到正文段落")


def test_absolute_spacing_becomes_a_multiple_of_the_font_size(tmp_path: Path):
    """21pt of spacing on 13.5pt text is 1.56 line-heights, not 266700."""
    # spcPts is in hundredths of a point: 2100 = 21pt.
    deck = _deck(tmp_path, spacing_xml=("a:spcPts", "2100"))

    spacing = _body_spacing(deck, tmp_path / "a")

    assert spacing == pytest.approx(21.0 / 13.5, rel=1e-3)


def test_absolute_spacing_never_reaches_the_renderer_as_raw_emu(tmp_path: Path):
    """The failure mode was silent: a plausible slide missing whole paragraphs."""
    deck = _deck(tmp_path, spacing_xml=("a:spcPts", "2100"))

    spacing = _body_spacing(deck, tmp_path / "b")

    # 21pt is 266700 EMU. Anything near that is the raw value leaking through.
    assert spacing < 5.0


def test_percentage_spacing_is_taken_as_the_multiple_it_already_is(tmp_path: Path):
    """spcPct is the common case and must not regress while fixing spcPts."""
    deck = _deck(tmp_path, spacing_xml=("a:spcPct", "150000"))  # 150%

    spacing = _body_spacing(deck, tmp_path / "c")

    assert spacing == pytest.approx(1.5, rel=1e-3)


def test_no_spacing_falls_back_to_the_default(tmp_path: Path):
    deck = _deck(tmp_path, spacing_xml=None)

    assert _body_spacing(deck, tmp_path / "d") == pytest.approx(1.2)
