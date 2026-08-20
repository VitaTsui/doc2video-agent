"""Flows read out of a deck, and the ones that must be left as pictures.

The gate this turns on: a connector in OOXML names the shapes it starts and
ends at, so the graph is a fact about the file. A slide whose arrows are loose
lines gives nothing here, and its picture stays its picture — the project only
rebuilds what it can read exactly (方案 §12、§20).
"""

from __future__ import annotations

from pathlib import Path

import pytest

from doc2video.schemas import BBox, DiagramEdge, DiagramFacts, DocumentPage, SlideElement


def _assets(path: Path) -> Path:
    """The parser writes page renders here; it does not create the directory."""
    path.mkdir(parents=True, exist_ok=True)
    return path


def _page(diagram: DiagramFacts | None) -> DocumentPage:
    return DocumentPage(
        index=1,
        width=1920,
        height=1080,
        diagram=diagram,
        elements=[
            SlideElement(id=name, text=name, bbox=BBox(x=0, y=0, w=100, h=50), label=name)
            for name in ("user", "gateway", "llm", "store")
        ],
    )


def test_the_flow_is_walked_in_the_order_its_arrows_go():
    """A flow is told in the order it flows, not in the order it was drawn."""
    facts = DiagramFacts(
        # Deliberately not in flow order: this is the order the shapes happen
        # to appear on the slide.
        nodes=["llm", "user", "store", "gateway"],
        edges=[
            DiagramEdge(source="user", target="gateway"),
            DiagramEdge(source="gateway", target="llm"),
            DiagramEdge(source="llm", target="store"),
        ],
    )
    assert facts.order() == ["user", "gateway", "llm", "store"]


def test_a_cycle_keeps_every_node():
    """A loop is a diagram too; dropping its nodes would drop the page."""
    facts = DiagramFacts(
        nodes=["a", "b", "c"],
        edges=[
            DiagramEdge(source="a", target="b"),
            DiagramEdge(source="b", target="c"),
            DiagramEdge(source="c", target="a"),
        ],
    )
    assert sorted(facts.order()) == ["a", "b", "c"]


def test_the_flow_is_told_to_the_writer_not_used_to_move_the_camera():
    """Where a declared flow is allowed to change the video, and where it is not.

    Re-pointing the camera along the arrows was tried and backed out: a shot is
    bound to the sentence that mentions it, so reordering the shots points the
    camera at a box the current sentence is not talking about. That breaks the
    one binding this project is built on (方案 §20) to fix a smaller problem.

    The flow belongs in front of the writer instead, where it can change the
    order things are *said* — and the shots then follow the sentences as they
    always did.
    """
    from doc2video.skills.director import DirectorSkill
    from doc2video.skills.narration import _flow_of

    page = _page(
        DiagramFacts(
            nodes=["llm", "user", "gateway"],
            edges=[
                DiagramEdge(source="user", target="gateway"),
                DiagramEdge(source="gateway", target="llm"),
            ],
        )
    )
    assert _flow_of(page) == "user → gateway → llm"
    assert _flow_of(_page(None)) == ""
    assert not hasattr(DirectorSkill, "_follow_the_arrows")


@pytest.fixture
def flow_deck(tmp_path: Path) -> Path:
    """A deck whose arrows declare what they join."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    boxes = []
    for name, x in (("用户", 1), ("网关", 4.5), ("大模型", 8), ("知识库", 11)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3), Inches(2), Inches(1)
        )
        shape.text_frame.text = name
        boxes.append(shape)
    for first, second in zip(boxes, boxes[1:], strict=False):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0), Inches(1), Inches(1)
        )
        line.begin_connect(first, 3)
        line.end_connect(second, 1)

    path = tmp_path / "flow.pptx"
    prs.save(path)
    return path


def test_a_declared_flow_is_read_out_of_the_file(flow_deck: Path, tmp_path: Path):
    """Read, not recognised: nothing here looks at the picture."""
    from doc2video.tools.parsers.ppt_parser import parse_ppt

    page = parse_ppt(flow_deck, _assets(tmp_path / "assets"), target_width=1920).pages[0]
    assert page.diagram is not None

    text = {element.id: element.text for element in page.elements}
    assert [text[node] for node in page.diagram.order()] == ["用户", "网关", "大模型", "知识库"]


def test_loose_arrows_leave_the_picture_a_picture(tmp_path: Path):
    """The other half of the gate, and the more important one.

    An architecture drawn with unattached lines cannot be read this way. The
    honest answer is nothing at all — a graph guessed from geometry would let
    the video explain a structure the slide never claimed.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Inches

    from doc2video.tools.parsers.ppt_parser import parse_ppt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for name, x in (("用户", 1), ("网关", 5)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3), Inches(2), Inches(1)
        )
        shape.text_frame.text = name
    # Lines that touch nothing: drawn between the boxes, attached to neither.
    for offset in (0, 1):
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(3 + offset), Inches(3.5), Inches(5), Inches(3.5)
        )

    path = tmp_path / "loose.pptx"
    prs.save(path)
    assert parse_ppt(path, _assets(tmp_path / "a"), target_width=1920).pages[0].diagram is None
