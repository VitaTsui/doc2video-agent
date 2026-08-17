"""Pattern fills and WordArt.

Both used to be dropped silently: a hatched shape came out transparent, and a
WordArt title came out as flat text — often invisible, because its colour lived
entirely in an outline or a gradient.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from doc2video.tools.slides.effects import text_effects
from doc2video.tools.slides.extract import extract_deck
from doc2video.tools.slides.pattern import mix, pattern_css
from doc2video.tools.slides.theme import Theme

BLUE = "#4472C4"
WHITE = "#FFFFFF"


# --------------------------------------------------------------------------
# pattern presets
# --------------------------------------------------------------------------


def test_upward_and_downward_diagonals_are_not_the_same_stripe():
    """CSS angles name the gradient axis; the stripes run across it.

    Reading them as the stripe direction mirrored every diagonal preset.
    """
    up = pattern_css("LIGHT_UPWARD_DIAGONAL", BLUE, WHITE)
    down = pattern_css("LIGHT_DOWNWARD_DIAGONAL", BLUE, WHITE)

    assert "135deg" in up and "45deg" in down
    assert up != down


def test_a_hatch_keeps_both_of_its_colours():
    css = pattern_css("CROSS", BLUE, WHITE)
    assert css.count("repeating-linear-gradient") == 2, "十字线需要两组线"
    assert BLUE in css and WHITE in css


def test_percent_shading_becomes_a_blend():
    assert pattern_css("PERCENT_50", "#000000", "#FFFFFF") == "#808080"
    assert pattern_css("PERCENT_25", "#000000", "#FFFFFF") == "#BFBFBF"


def test_an_unknown_preset_still_produces_a_fill():
    """Never nothing: a transparent shape is further from the slide than a blend."""
    assert pattern_css("SOMETHING_POWERPOINT_ADDED_LATER", BLUE, WHITE)
    assert pattern_css("MIXED", BLUE, WHITE) == mix(BLUE, WHITE, 0.5)


@pytest.mark.parametrize(
    "preset",
    [member for member in dir(MSO_PATTERN_TYPE) if member.isupper()],
)
def test_every_preset_powerpoint_offers_resolves(preset: str):
    assert pattern_css(preset, BLUE, WHITE)


def test_patterned_shape_reaches_the_render_model(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
    box.text_frame.text = "图案"
    box.fill.patterned()
    box.fill.pattern = MSO_PATTERN_TYPE.LIGHT_UPWARD_DIAGONAL
    box.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    box.fill.back_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    path = tmp_path / "pattern.pptx"
    prs.save(path)

    deck = extract_deck(path, tmp_path / "assets", target_width=1920)
    fills = [s.style.fill for s in deck.slides[0].shapes if s.style.fill]

    assert fills, "图案填充不应被丢成透明"
    assert "repeating-linear-gradient" in fills[0]


# --------------------------------------------------------------------------
# WordArt
# --------------------------------------------------------------------------


def _wordart_deck(tmp_path: Path, *, gradient=True, outline=True, shadow=True) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(2))
    box.text_frame.text = "艺术字"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(54)
    rpr = run._r.find(qn("a:rPr"))

    if gradient:
        grad = etree.SubElement(rpr, qn("a:gradFill"))
        stops = etree.SubElement(grad, qn("a:gsLst"))
        for pos, color in (("0", "4472C4"), ("100000", "ED7D31")):
            stop = etree.SubElement(stops, qn("a:gs"))
            stop.set("pos", pos)
            etree.SubElement(stop, qn("a:srgbClr")).set("val", color)
        etree.SubElement(grad, qn("a:lin")).set("ang", "2700000")
    if outline:
        line = etree.SubElement(rpr, qn("a:ln"))
        line.set("w", "12700")
        fill = etree.SubElement(line, qn("a:solidFill"))
        etree.SubElement(fill, qn("a:srgbClr")).set("val", "1F2430")
    if shadow:
        effects = etree.SubElement(rpr, qn("a:effectLst"))
        shdw = etree.SubElement(effects, qn("a:outerShdw"))
        shdw.set("blurRad", "50800")
        shdw.set("dist", "38100")
        shdw.set("dir", "2700000")
        color = etree.SubElement(shdw, qn("a:srgbClr"))
        color.set("val", "000000")
        etree.SubElement(color, qn("a:alpha")).set("val", "40000")

    path = tmp_path / "wordart.pptx"
    prs.save(path)
    return path


def _first_run(deck):
    for shape in deck.slides[0].shapes:
        if shape.text and shape.text.paragraphs:
            return shape.text.paragraphs[0].runs[0]
    raise AssertionError("没有提取到文字")


def test_wordart_outline_gradient_and_shadow_all_survive(tmp_path: Path):
    deck = extract_deck(
        _wordart_deck(tmp_path), tmp_path / "assets", target_width=1920
    )
    effects = _first_run(deck).effects

    assert effects is not None
    assert effects.outline_color == "#1F2430"
    assert effects.outline_width_px > 0
    assert effects.gradient is not None
    assert "#4472C4" in effects.gradient and "#ED7D31" in effects.gradient
    assert effects.shadow is not None


def test_reading_the_colour_must_not_eat_the_gradient(tmp_path: Path):
    """python-pptx's ColorFormat calls get_or_change_to_solidFill().

    Asking a run for `font.color` therefore *replaces* its gradient fill with an
    empty solid one, so the effects have to be read first. This pins the order.
    """
    path = _wordart_deck(tmp_path, outline=False, shadow=False)
    prs = Presentation(path)
    run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]

    _ = run.font.color.type  # the destructive read
    assert text_effects(run, Theme(), 1.0) is None, "本测试假设读取颜色会破坏渐变"

    # The extractor must not hit that order, so the gradient still arrives.
    deck = extract_deck(path, tmp_path / "assets", target_width=1920)
    assert _first_run(deck).effects.gradient is not None


def test_shadow_alpha_becomes_a_translucent_colour(tmp_path: Path):
    deck = extract_deck(
        _wordart_deck(tmp_path, gradient=False, outline=False),
        tmp_path / "assets",
        target_width=1920,
    )
    assert "rgba(0, 0, 0, 0.40)" in _first_run(deck).effects.shadow


def test_a_width_with_no_colour_draws_no_outline(tmp_path: Path):
    """An `a:ln` with no fill of its own would mean guessing at the palette."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    box.text_frame.text = "描边无色"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(40)
    line = etree.SubElement(run._r.find(qn("a:rPr")), qn("a:ln"))
    line.set("w", "25400")
    path = tmp_path / "noclr.pptx"
    prs.save(path)

    deck = extract_deck(path, tmp_path / "assets", target_width=1920)
    assert _first_run(deck).effects is None


def test_plain_text_carries_no_effects(demo_pptx: Path, tmp_path: Path):
    deck = extract_deck(demo_pptx, tmp_path / "assets", target_width=1920)
    runs = [
        run
        for slide in deck.slides
        for shape in slide.shapes
        if shape.text
        for para in shape.text.paragraphs
        for run in para.runs
    ]
    assert runs
    assert all(run.effects is None for run in runs)
