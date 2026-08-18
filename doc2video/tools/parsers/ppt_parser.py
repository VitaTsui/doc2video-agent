"""PPT / PPTX parser built on python-pptx (+ LibreOffice when available).

Shape geometry comes from the OOXML package, which is exact. Slide *rendering*
has two paths: LibreOffice → PDF → PyMuPDF (pixel-accurate, preserves theme and
native animations' end state), or a Pillow rasterizer from the same geometry
(always available, plain styling).
"""

from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from ...core.config import Settings, which
from ...core.errors import UnsupportedSource
from ...core.ids import element_id
from ...core.logging import get_logger
from ...schemas import BBox, DocumentModel, DocumentPage, ElementKind, SlideElement
from .slide_raster import rasterize_page

log = get_logger(__name__)

SOFFICE_TIMEOUT = 180


def parse_ppt(
    path: Path,
    assets_dir: Path,
    *,
    target_width: int = 1920,
    settings: Settings | None = None,
) -> DocumentModel:
    source = path
    if path.suffix.lower() == ".ppt":
        source = _convert_legacy_ppt(path, assets_dir)

    prs = Presentation(source)
    slide_w = Emu(prs.slide_width or 9144000)
    slide_h = Emu(prs.slide_height or 6858000)
    scale = target_width / slide_w.inches / 96 if slide_w.inches else 1.0
    px_w = target_width
    px_h = int(round(slide_h.inches * 96 * scale))

    pages: list[DocumentPage] = []
    for index, slide in enumerate(prs.slides, start=1):
        elements = _extract_elements(slide, index, slide_w, slide_h, px_w, px_h, assets_dir)
        pages.append(
            DocumentPage(
                index=index,
                title=_slide_title(slide, elements),
                elements=elements,
                speaker_notes=_notes(slide),
                width=float(px_w),
                height=float(px_h),
            )
        )

    _render_pages(source, pages, assets_dir, target_width=target_width, settings=settings)

    log.info("解析 PPT 完成：%s，共 %d 页", path.name, len(pages))
    return DocumentModel(
        title=path.stem,
        pages=pages,
        presentation_order=[p.index for p in pages],
    )


# --------------------------------------------------------------------------
# elements
# --------------------------------------------------------------------------


def _extract_elements(
    slide, page_index: int, slide_w, slide_h, px_w: int, px_h: int, assets_dir: Path
) -> list[SlideElement]:
    elements: list[SlideElement] = []
    font_sizes: dict[str, float] = {}
    seq = 0

    for shape in slide.shapes:
        bbox = _shape_bbox(shape, slide_w, slide_h, px_w, px_h)
        if bbox is None:
            continue
        seq += 1

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            asset = _save_picture(shape, page_index, seq, assets_dir)
            elements.append(
                SlideElement(
                    id=element_id(page_index, seq, "image"),
                    kind=ElementKind.IMAGE,
                    bbox=bbox,
                    label=f"image_{seq}",
                    importance=0.6,
                    asset_path=asset,
                )
            )
            continue

        if shape.has_chart:
            elements.append(
                SlideElement(
                    id=element_id(page_index, seq, "chart"),
                    kind=ElementKind.CHART,
                    text=_chart_text(shape),
                    bbox=bbox,
                    label=f"chart_{seq}",
                    importance=0.85,
                )
            )
            continue

        if shape.has_table:
            elements.append(
                SlideElement(
                    id=element_id(page_index, seq, "table"),
                    kind=ElementKind.TABLE,
                    text=_table_text(shape),
                    bbox=bbox,
                    label=f"table_{seq}",
                    importance=0.75,
                )
            )
            continue

        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            kind, importance, level = _classify_text_shape(shape)
            elements.append(
                SlideElement(
                    id=element_id(page_index, seq, text),
                    kind=kind,
                    text=text,
                    bbox=bbox,
                    level=level,
                    label=text.split("\n")[0][:24],
                    importance=importance,
                )
            )
            font_sizes[elements[-1].id] = _max_font_size(shape)

    _promote_title(elements, font_sizes)
    return elements


def _classify_text_shape(shape) -> tuple[ElementKind, float, int]:
    # A title placeholder (idx 0) is authoritative when the deck uses layouts.
    if _placeholder_index(shape) == 0:
        return ElementKind.TITLE, 0.9, 0
    paragraphs = shape.text_frame.paragraphs
    level = max((p.level for p in paragraphs), default=0)
    if len(paragraphs) > 1 or level > 0:
        return ElementKind.BULLET, 0.6, level
    return ElementKind.PARAGRAPH, 0.5, level


def _placeholder_index(shape) -> int | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return shape.placeholder_format.idx
    except ValueError:
        return None


def _max_font_size(shape) -> float:
    """Largest explicit run size on the shape, in points; 0 when inherited."""
    largest = 0.0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                largest = max(largest, run.font.size.pt)
    return largest


def _promote_title(elements: list[SlideElement], font_sizes: dict[str, float]) -> None:
    """Decks built from plain text boxes have no title placeholder.

    Fall back to the largest explicitly-sized text on the slide — the same
    signal the PDF parser uses — so every page still has an addressable title.
    """
    if any(el.kind is ElementKind.TITLE for el in elements):
        return
    candidates = [
        (font_sizes.get(el.id, 0.0), el)
        for el in elements
        if el.kind in (ElementKind.PARAGRAPH, ElementKind.BULLET) and el.text
    ]
    if not candidates:
        return
    largest_size, candidate = max(candidates, key=lambda item: (item[0], -item[1].bbox.y))
    if largest_size <= 0:
        # No explicit sizes at all: the topmost text box is the best guess.
        candidate = min((el for _, el in candidates), key=lambda el: el.bbox.y)
    candidate.kind = ElementKind.TITLE
    candidate.importance = 0.9


def _shape_bbox(shape, slide_w, slide_h, px_w: int, px_h: int) -> BBox | None:
    if shape.left is None or shape.top is None or not shape.width or not shape.height:
        return None
    x = float(shape.left) / float(slide_w) * px_w
    y = float(shape.top) / float(slide_h) * px_h
    w = float(shape.width) / float(slide_w) * px_w
    h = float(shape.height) / float(slide_h) * px_h
    if w < 8 or h < 8:
        return None
    return BBox(x=x, y=y, w=w, h=h)


def _save_picture(shape, page_index: int, seq: int, assets_dir: Path) -> str | None:
    try:
        image = shape.image
    except (AttributeError, ValueError):
        return None
    name = f"p{page_index:03d}_img{seq:02d}.{image.ext}"
    (assets_dir / name).write_bytes(image.blob)
    return str(assets_dir / name)


def _chart_text(shape) -> str:
    """Describe a chart with its actual numbers.

    The narration skill can only avoid reading the slide aloud if it is told
    what the chart *says*; a list of category names is not enough.
    """
    from ..slides import chart_of, describe_chart

    try:
        chart = chart_of(shape)
    except Exception:  # malformed chart part
        return "图表"
    if chart is None:
        return "图表"
    return describe_chart(chart)


def _table_text(shape) -> str:
    rows = []
    for row in shape.table.rows:
        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(rows[:12])


def _slide_title(slide, elements: list[SlideElement]) -> str:
    try:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except (AttributeError, ValueError):
        pass
    for element in elements:
        if element.kind is ElementKind.TITLE and element.text:
            return element.text
    return ""


def _notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    return frame.text.strip() if frame is not None else ""


# --------------------------------------------------------------------------
# rendering
# --------------------------------------------------------------------------


def _render_pages(
    source: Path,
    pages: list[DocumentPage],
    assets_dir: Path,
    *,
    target_width: int,
    settings: Settings | None = None,
) -> None:
    """Rasterize slides, in descending order of fidelity.

    1. LibreOffice — PowerPoint's own layout engine, closest to the original.
    2. Chromium — reuses Remotion's browser; keeps theme colours, fonts, fills,
       gradients and rotation, without a 1GB office suite.
    3. Pillow — geometry only, but always available.
    """
    if _render_with_libreoffice(source, pages, assets_dir, target_width=target_width):
        return

    if _render_with_chromium(
        source, pages, assets_dir, target_width=target_width, settings=settings
    ):
        return

    log.warning("LibreOffice 与 Chromium 均不可用，使用内置栅格化器（样式为简版）")
    for page in pages:
        name = f"page_{page.index:03d}.png"
        rasterize_page(page, assets_dir / name)
        page.image_path = f"assets/{name}"


def _render_with_chromium(
    source: Path,
    pages: list[DocumentPage],
    assets_dir: Path,
    *,
    target_width: int,
    settings: Settings | None = None,
) -> bool:
    from ..slides import ChromiumSlideRenderer, extract_deck

    renderer = ChromiumSlideRenderer(settings)
    if not renderer.available():
        log.info("Chromium 幻灯片渲染不可用：%s", renderer.unavailable_reason())
        return False

    try:
        deck = extract_deck(source, assets_dir, target_width=target_width)
        paths = renderer.render(deck, assets_dir)
    except Exception as exc:
        log.warning("Chromium 幻灯片渲染失败，回退到内置栅格化器：%s", exc)
        return False

    by_index = {slide.index: path for slide, path in zip(deck.slides, paths, strict=False)}
    for page in pages:
        path = by_index.get(page.index)
        if path is None:
            return False
        page.image_path = path
        # The deck render is authoritative for geometry; keep element boxes in step.
        _rescale_page(page, float(deck.width), float(deck.height))

    log.info("使用 Chromium 渲染幻灯片（%d 页）", len(paths))
    return True


def _render_with_libreoffice(
    source: Path, pages: list[DocumentPage], assets_dir: Path, *, target_width: int
) -> bool:
    if which("soffice") is None:
        return False
    try:
        import fitz
    except ImportError:  # pragma: no cover - PyMuPDF is a hard dependency
        return False

    with tempfile.TemporaryDirectory() as tmp:
        try:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, str(source)],
                check=True,
                capture_output=True,
                timeout=SOFFICE_TIMEOUT,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
            log.warning("LibreOffice 转 PDF 失败，回退到内置栅格化器：%s", exc)
            return False

        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            return False

        doc = fitz.open(pdfs[0])
        try:
            for page, pdf_page in zip(pages, doc, strict=False):
                zoom = target_width / pdf_page.rect.width if pdf_page.rect.width else 1.0
                pixmap = pdf_page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
                name = f"page_{page.index:03d}.png"
                pixmap.save(assets_dir / name)
                page.image_path = f"assets/{name}"
                # Trust the actual render for page geometry; rescale element boxes.
                _rescale_page(page, float(pixmap.width), float(pixmap.height))
        finally:
            doc.close()
    return True


def _rescale_page(page: DocumentPage, new_w: float, new_h: float) -> None:
    if not page.width or not page.height:
        page.width, page.height = new_w, new_h
        return
    sx, sy = new_w / page.width, new_h / page.height
    if abs(sx - 1) < 1e-6 and abs(sy - 1) < 1e-6:
        return
    for element in page.elements:
        element.bbox = BBox(
            x=element.bbox.x * sx,
            y=element.bbox.y * sy,
            w=element.bbox.w * sx,
            h=element.bbox.h * sy,
        )
    page.width, page.height = new_w, new_h


def _convert_legacy_ppt(path: Path, assets_dir: Path) -> Path:
    """Legacy binary .ppt is not readable by python-pptx; convert it first."""
    if which("soffice") is None:
        raise UnsupportedSource(
            "旧版 .ppt 需要 LibreOffice 才能解析，请安装 LibreOffice 或先另存为 .pptx",
            detail={"missing": "soffice"},
        )
    out_dir = assets_dir / "_converted"
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pptx", "--outdir", str(out_dir), str(path)],
        check=True,
        capture_output=True,
        timeout=SOFFICE_TIMEOUT,
    )
    converted = list(out_dir.glob("*.pptx"))
    if not converted:
        raise UnsupportedSource("LibreOffice 未能将 .ppt 转换为 .pptx")
    return converted[0]
