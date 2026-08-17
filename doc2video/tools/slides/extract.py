"""python-pptx → SlideDeck.

Reads the styling the semantic parser deliberately ignores: theme colours,
fonts, fills, gradients, line work, rotation, group transforms, tables.

Placeholder inheritance is resolved too (see ``inherit.py``): a body
placeholder's size, colour, alignment and bullet usually live in the layout or
master rather than on the slide, so reading only the slide yields the text
without its appearance. ``DEFAULT_SIZES`` now applies only when even the master
says nothing.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Length

from ...core.logging import get_logger
from .chart_xml import PlotGroups, RawSeries, read_plot_groups
from .effects import text_effects
from .inherit import LevelDefaults, StyleResolver
from .model import (
    Align,
    Box,
    ChartData,
    ChartKind,
    ChartSeries,
    Geometry,
    Paragraph,
    Run,
    ShapeKind,
    ShapeStyle,
    Slide,
    SlideDeck,
    SlideShape,
    TableCell,
    TableData,
    TextBody,
    VAnchor,
)
from .pattern import pattern_css
from .table_style import TableStyle, TableStyles, load_table_styles
from .theme import Theme, load_theme

log = get_logger(__name__)

EMU_PER_INCH = 914400
EMU_PER_POINT = 12700

# Fallback sizes for runs whose size is inherited from the layout.
DEFAULT_SIZES = {
    PP_PLACEHOLDER.TITLE: 44.0,
    PP_PLACEHOLDER.CENTER_TITLE: 44.0,
    PP_PLACEHOLDER.SUBTITLE: 32.0,
    PP_PLACEHOLDER.BODY: 18.0,
    PP_PLACEHOLDER.OBJECT: 18.0,
}
DEFAULT_SIZE = 18.0
DEFAULT_LINE_SPACING = 1.2

ALIGN_MAP = {
    PP_ALIGN.CENTER: Align.CENTER,
    PP_ALIGN.RIGHT: Align.RIGHT,
    PP_ALIGN.JUSTIFY: Align.JUSTIFY,
    PP_ALIGN.JUSTIFY_LOW: Align.JUSTIFY,
    PP_ALIGN.LEFT: Align.LEFT,
}

ANCHOR_MAP = {
    MSO_ANCHOR.MIDDLE: VAnchor.MIDDLE,
    MSO_ANCHOR.BOTTOM: VAnchor.BOTTOM,
    MSO_ANCHOR.TOP: VAnchor.TOP,
}

GEOMETRY_MAP = {
    MSO_SHAPE.ROUNDED_RECTANGLE: Geometry.ROUND_RECT,
    MSO_SHAPE.OVAL: Geometry.ELLIPSE,
    MSO_SHAPE.RECTANGLE: Geometry.RECT,
}


def extract_deck(
    pptx_path: Path, assets_dir: Path, *, target_width: int
) -> SlideDeck:
    """Build the render model for a whole presentation."""
    prs = Presentation(pptx_path)
    theme = load_theme(pptx_path)

    slide_w_emu = int(prs.slide_width or 9144000)
    slide_h_emu = int(prs.slide_height or 6858000)
    emu_to_px = target_width / slide_w_emu
    px_height = int(round(slide_h_emu * emu_to_px))
    pt_to_px = target_width / (slide_w_emu / EMU_PER_INCH) / 72

    ctx = _Context(
        theme=theme,
        emu_to_px=emu_to_px,
        assets_dir=assets_dir,
        styles=StyleResolver(prs, theme),
        table_styles=load_table_styles(pptx_path, theme),
    )

    slides: list[Slide] = []
    for index, slide in enumerate(prs.slides, start=1):
        ctx.page_index = index
        ctx.seq = 0
        shapes: list[SlideShape] = []
        for shape in slide.shapes:
            shapes.extend(ctx.convert(shape))
        slides.append(
            Slide(
                index=index,
                background=_slide_background(slide, theme),
                shapes=shapes,
            )
        )

    return SlideDeck(
        width=target_width, height=px_height, pt_to_px=pt_to_px, slides=slides
    )


def chart_of(shape, theme: Theme | None = None) -> ChartData | None:
    """Extract just a chart, for callers that only need its data (not styling)."""
    ctx = _Context(theme=theme or Theme(), emu_to_px=1.0, assets_dir=Path("."))
    return ctx._chart(shape)


class _Context:
    """Carries the per-deck constants through the recursive shape walk."""

    def __init__(
        self,
        *,
        theme: Theme,
        emu_to_px: float,
        assets_dir: Path,
        styles: StyleResolver | None = None,
        table_styles: TableStyles | None = None,
    ) -> None:
        self.theme = theme
        self.emu_to_px = emu_to_px
        self.assets_dir = assets_dir
        self.styles = styles
        self.table_styles = table_styles
        self.page_index = 0
        self.seq = 0

    # -- shape dispatch --------------------------------------------------
    def convert(self, shape, transform: _GroupTransform | None = None) -> list[SlideShape]:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return self._convert_group(shape, transform)

        box = self._box(shape, transform)
        if box is None:
            return []

        style = self._style(shape, box)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = self._save_picture(shape)
            if image is None:
                return []
            return [SlideShape(kind=ShapeKind.PICTURE, box=box, style=style, image=image)]

        if shape.has_table:
            return [
                SlideShape(
                    kind=ShapeKind.TABLE, box=box, style=style, table=self._table(shape, box)
                )
            ]

        if shape.has_chart:
            chart = self._chart(shape)
            if chart is not None:
                return [SlideShape(kind=ShapeKind.CHART, box=box, style=style, chart=chart)]
            return []

        text = self._text(shape) if shape.has_text_frame else None
        has_visible_fill = style.fill is not None or style.line_color is not None
        if text is None and not has_visible_fill:
            # Nothing to draw: an empty placeholder or a chart we cannot render.
            return []

        kind = ShapeKind.TEXT if shape.has_text_frame else ShapeKind.AUTO
        return [SlideShape(kind=kind, box=box, style=style, text=text)]

    def _convert_group(self, group, parent: _GroupTransform | None) -> list[SlideShape]:
        """Flatten a group, mapping child coordinates through its transform.

        Grouped shapes store coordinates in the group's own child space, so
        without applying ``chOff``/``chExt`` every icon in a diagram lands in the
        wrong place — usually stacked in a corner.
        """
        transform = _GroupTransform.of(group, parent)
        out: list[SlideShape] = []
        for child in group.shapes:
            out.extend(self.convert(child, transform))
        return out

    # -- geometry ---------------------------------------------------------
    def _box(self, shape, transform: _GroupTransform | None) -> Box | None:
        left, top = shape.left, shape.top
        width, height = shape.width, shape.height
        if left is None or top is None or not width or not height:
            inherited = _inherited_geometry(shape)
            if inherited is None:
                return None
            left, top, width, height = inherited

        if transform is not None:
            left, top, width, height = transform.apply(left, top, width, height)

        box = Box(
            x=left * self.emu_to_px,
            y=top * self.emu_to_px,
            w=width * self.emu_to_px,
            h=height * self.emu_to_px,
        )
        return box if box.w >= 1 and box.h >= 1 else None

    def _style(self, shape, box: Box) -> ShapeStyle:
        style = ShapeStyle(
            fill=self._fill(shape),
            line_color=self._line_color(shape),
            line_width_px=self._line_width(shape),
            geometry=self._geometry(shape),
            rotation=float(getattr(shape, "rotation", 0.0) or 0.0),
        )
        if style.geometry is Geometry.ROUND_RECT:
            style.corner_radius_px = _corner_radius(shape, min(box.w, box.h))
        return style

    def _geometry(self, shape) -> Geometry:
        try:
            auto_shape = shape.auto_shape_type
        except (AttributeError, ValueError):
            return Geometry.RECT
        if auto_shape is None:
            return Geometry.RECT
        return GEOMETRY_MAP.get(auto_shape, Geometry.OTHER)

    def _fill(self, shape) -> str | None:
        try:
            return self._fill_value(shape.fill)
        except (AttributeError, ValueError, NotImplementedError):
            return None

    def _fill_value(self, fill) -> str | None:
        """Any fill python-pptx can name, as a CSS ``background`` value."""
        try:
            fill_type = fill.type
        except (AttributeError, ValueError, NotImplementedError):
            return None
        if fill_type is None:
            return None

        type_name = getattr(fill_type, "name", str(fill_type))
        if type_name == "SOLID":
            return self.theme.resolve(fill.fore_color)
        if type_name == "GRADIENT":
            return self._gradient(fill)
        if type_name == "PATTERNED":
            return self._pattern(fill)
        # BACKGROUND means "no fill"; picture fills are staged as images instead.
        return None

    def _pattern(self, fill) -> str | None:
        """A hatch preset as CSS.

        Dropping the fill would make a hatched shape transparent, which reads as
        a different shape rather than an approximate one — so an unmapped preset
        still comes back as a blend of its two colours.
        """
        try:
            preset = getattr(fill.pattern, "name", "") or ""
            fore = self.theme.resolve(fill.fore_color)
            back = self.theme.resolve(fill.back_color)
        except (AttributeError, ValueError, NotImplementedError):
            return None
        if not fore and not back:
            return None
        return pattern_css(preset, fore or "#000000", back or "#FFFFFF")

    def _gradient(self, fill) -> str | None:
        try:
            stops = list(fill.gradient_stops)
        except (AttributeError, ValueError, NotImplementedError):
            return None
        colors = []
        for stop in stops:
            color = self.theme.resolve(stop.color)
            if color:
                colors.append(f"{color} {round(stop.position * 100)}%")
        if len(colors) < 2:
            return colors[0].split(" ")[0] if colors else None
        try:
            angle = float(fill.gradient_angle or 0.0)
        except (AttributeError, ValueError, NotImplementedError, TypeError):
            angle = 0.0
        # OOXML measures clockwise from the positive x-axis; CSS from "up".
        return f"linear-gradient({angle + 90:.0f}deg, {', '.join(colors)})"

    def _line_color(self, shape) -> str | None:
        try:
            line = shape.line
            if line.fill.type is None:
                return None
            type_name = getattr(line.fill.type, "name", "")
            if type_name != "SOLID":
                return None
            return self.theme.resolve(line.color)
        except (AttributeError, ValueError, NotImplementedError):
            return None

    def _line_width(self, shape) -> float:
        try:
            width = shape.line.width
        except (AttributeError, ValueError, NotImplementedError):
            return 0.0
        if not width:
            return 0.0
        return float(width) * self.emu_to_px

    # -- text -------------------------------------------------------------
    def _text(self, shape) -> TextBody | None:
        frame = shape.text_frame
        if not frame.text.strip():
            return None

        scale = self.styles.font_scale(shape) if self.styles else 1.0
        paragraphs = [self._paragraph(p, shape, scale) for p in frame.paragraphs]
        paragraphs = [p for p in paragraphs if p.runs]
        if not paragraphs:
            return None

        return TextBody(
            paragraphs=paragraphs,
            v_anchor=ANCHOR_MAP.get(frame.vertical_anchor, VAnchor.TOP),
            wrap=frame.word_wrap is not False,
            margin_left_pt=_emu_to_pt(frame.margin_left, 7.2),
            margin_right_pt=_emu_to_pt(frame.margin_right, 7.2),
            margin_top_pt=_emu_to_pt(frame.margin_top, 3.6),
            margin_bottom_pt=_emu_to_pt(frame.margin_bottom, 3.6),
            default_size_pt=_default_size(shape) * scale,
        )

    def _paragraph(self, paragraph, shape=None, scale: float = 1.0) -> Paragraph:
        level = int(paragraph.level or 0)
        # What the layout/master says for this placeholder at this level; every
        # value here is a fallback the slide itself may override.
        inherited = (
            self.styles.defaults_for(shape, level)
            if self.styles is not None and shape is not None
            else LevelDefaults()
        )

        runs = []
        for run in paragraph.runs:
            if not run.text:
                continue
            # Read WordArt first: python-pptx's ColorFormat calls
            # get_or_change_to_solidFill(), so merely *asking* for run.font.color
            # replaces a gradient text fill with an empty solid one.
            effects = text_effects(run, self.theme, self.emu_to_px)
            explicit_size = run.font.size.pt if run.font.size is not None else None
            size_pt = explicit_size if explicit_size is not None else inherited.size_pt
            runs.append(
                Run(
                    text=run.text,
                    bold=_first_set(run.font.bold, inherited.bold, False),
                    italic=_first_set(run.font.italic, inherited.italic, False),
                    underline=_first_set(run.font.underline, inherited.underline, False),
                    size_pt=size_pt * scale if size_pt is not None else None,
                    color=self.theme.resolve(run.font.color) or inherited.color,
                    font=run.font.name or inherited.font,
                    effects=effects,
                )
            )
        bullet = _bullet_of(paragraph, inherited=inherited)
        if bullet and runs:
            # Some decks type the bullet into the text as well; rendering both
            # would double it.
            runs[0].text = _strip_leading_bullet(runs[0].text)

        explicit_align = ALIGN_MAP.get(paragraph.alignment) if paragraph.alignment else None
        return Paragraph(
            runs=runs,
            align=explicit_align or Align(inherited.align or "left"),
            level=level,
            bullet=bullet,
            line_spacing=_line_spacing(paragraph, runs),
            space_before_pt=_length_to_pt(paragraph.space_before),
            space_after_pt=_length_to_pt(paragraph.space_after),
        )

    # -- other content ------------------------------------------------------
    def _table(self, shape, box: Box) -> TableData:
        table = shape.table
        rows: list[list[TableCell]] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                paragraphs = cell.text_frame.paragraphs
                first = paragraphs[0] if paragraphs else None
                bold = bool(first.runs and first.runs[0].font.bold) if first else False
                cells.append(
                    TableCell(
                        text=cell.text.strip(),
                        bold=bold,
                        align=ALIGN_MAP.get(first.alignment, Align.LEFT) if first else Align.LEFT,
                        fill=self._cell_fill(cell),
                    )
                )
            rows.append(cells)

        total_w = sum(int(c.width or 0) for c in table.columns) or 1
        total_h = sum(int(r.height or 0) for r in table.rows) or 1

        # PowerPoint always applies a table style. Read the real definition when
        # the deck ships one; otherwise approximate its out-of-the-box default
        # (accent header over banded rows), which is still far closer than
        # bare gridlines.
        style = self.table_styles.get(_table_style_id(table)) if self.table_styles else None
        accent = self.theme.slot("accent1")
        banded = _has_header(table)

        return TableData(
            rows=rows,
            col_widths=[int(c.width or 0) / total_w * box.w for c in table.columns],
            row_heights=[int(r.height or 0) / total_h * box.h for r in table.rows],
            header_fill=_style_value(style, "header_fill", accent) if banded else None,
            band_fill=_style_value(style, "band_fill", self.theme.slot("accent1", 0.8)),
            border_color=_style_value(style, "border_color", "#FFFFFF"),
            header_text=_style_value(style, "header_text", "#FFFFFF"),
        )

    def _chart(self, shape) -> ChartData | None:
        """Read a chart's data and formatting, faithful to the deck."""
        try:
            chart = shape.chart
        except (AttributeError, ValueError):
            return None

        groups = read_plot_groups(chart, self.theme)
        # The declared type names the whole chart; the plot groups know better
        # when it mixes types, and cover the 3D and cone/pyramid variants that
        # have no XL_CHART_TYPE entry of their own.
        kind = CHART_KIND_MAP.get(_chart_type_name(chart), ChartKind.OTHER)
        if kind is ChartKind.OTHER:
            kind = groups.kind_at(0) or ChartKind.OTHER
        categories = _chart_categories(chart) or groups.categories
        series = self._chart_series(chart, groups, kind)
        if not series:
            return None

        primary = groups.primary_scale
        secondary = groups.secondary_scale
        is_round = kind in (ChartKind.PIE, ChartKind.DOUGHNUT)
        return ChartData(
            kind=kind,
            title=_chart_title(chart),
            categories=categories,
            series=series,
            point_colors=self._point_colors(series, len(categories)) if is_round else [],
            # A legend for a single cartesian series is noise — the title already
            # names it. A pie is the opposite: its slices are only identifiable
            # through the legend, so honour the deck's setting there.
            has_legend=bool(getattr(chart, "has_legend", False))
            and (len(series) > 1 or is_round),
            legend_position=_legend_position(chart),
            has_data_labels=_has_data_labels(chart),
            gridlines=_has_gridlines(chart),
            # Once the plot area has parsed it is the authority: python-pptx
            # answers `value_axis` with one of the two axes on a combo chart,
            # and reading bars against the percentage axis pins every bar to
            # the ceiling. Its value is only a fallback for unparsable XML.
            y_min=primary.minimum if primary else _axis_scale(chart, "minimum_scale"),
            y_max=primary.maximum if primary else _axis_scale(chart, "maximum_scale"),
            y2_min=secondary.minimum if secondary else None,
            y2_max=secondary.maximum if secondary else None,
            y2_visible=secondary.visible if secondary else True,
            three_d=groups.three_d or _is_three_d(_chart_type_name(chart)),
        )

    def _chart_series(
        self, chart, groups: PlotGroups, chart_kind: ChartKind
    ) -> list[ChartSeries]:
        try:
            raw_series = list(chart.series)
        except (AttributeError, ValueError):
            # python-pptx opens nine plot tags; a 3D or of-pie chart raises here.
            raw_series = []

        series: list[ChartSeries] = []
        for index in range(max(len(raw_series), len(groups.series))):
            item = raw_series[index] if index < len(raw_series) else None
            from_xml = groups.series[index] if index < len(groups.series) else None
            own_kind = groups.kind_at(index)
            series.append(
                ChartSeries(
                    name=self._series_name(item, from_xml, index),
                    values=self._series_values(item, from_xml),
                    color=self._series_color(item, from_xml, index),
                    # Only carry a kind when it differs; every series naming the
                    # chart's own kind would just be noise in the model.
                    kind=own_kind if own_kind not in (None, chart_kind) else None,
                    secondary_axis=groups.secondary_at(index),
                )
            )
        return [s for s in series if s.values]

    @staticmethod
    def _series_name(item, from_xml: RawSeries | None, index: int) -> str:
        name = str(getattr(item, "name", "") or "") if item is not None else ""
        return name or (from_xml.name if from_xml else "") or f"系列 {index + 1}"

    @staticmethod
    def _series_values(item, from_xml: RawSeries | None) -> list[float | None]:
        if item is not None:
            try:
                values = [None if v is None else float(v) for v in item.values]
            except (AttributeError, TypeError, ValueError):
                values = []
            if values:
                return values
        return list(from_xml.values) if from_xml else []

    def _point_colors(self, series: list[ChartSeries], category_count: int) -> list[str]:
        """One colour per slice, the way PowerPoint assigns them.

        A pie's identity is the category, so deriving slice colours from the
        first series' colour leaves near-identical neighbours. PowerPoint walks
        the theme accents instead, which is both faithful and legible.
        """
        count = max(category_count, len(series[0].values) if series else 0)
        return [self.theme.slot(f"accent{i % 6 + 1}") for i in range(count)]

    def _series_color(self, item, from_xml: RawSeries | None, index: int) -> str:
        """Explicit series fill, else the theme accent PowerPoint would use.

        PowerPoint cycles accent1..accent6 by series index, so reproducing a deck
        means cycling too — even though a chart designed from scratch should not.
        """
        try:
            fill = item.format.fill if item is not None else None
            if fill is not None and getattr(fill.type, "name", "") == "SOLID":
                explicit = self.theme.resolve(fill.fore_color)
                if explicit:
                    return explicit
        except (AttributeError, ValueError, NotImplementedError):
            pass
        if from_xml is not None and from_xml.color:
            return from_xml.color
        return self.theme.slot(f"accent{index % 6 + 1}")

    def _cell_fill(self, cell) -> str | None:
        try:
            return self._fill_value(cell.fill)
        except (AttributeError, ValueError, NotImplementedError):
            return None

    def _save_picture(self, shape) -> str | None:
        try:
            image = shape.image
        except (AttributeError, ValueError):
            return None
        self.seq += 1
        name = f"slide{self.page_index:03d}_pic{self.seq:02d}.{image.ext}"
        (self.assets_dir / name).write_bytes(image.blob)
        return name


class _GroupTransform:
    """Maps a group's child coordinate space onto the slide."""

    def __init__(self, off_x, off_y, ch_off_x, ch_off_y, scale_x, scale_y, parent=None) -> None:
        self.off_x, self.off_y = off_x, off_y
        self.ch_off_x, self.ch_off_y = ch_off_x, ch_off_y
        self.scale_x, self.scale_y = scale_x, scale_y
        self.parent = parent

    @classmethod
    def of(cls, group, parent: _GroupTransform | None) -> _GroupTransform | None:
        try:
            xfrm = group._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            ch_off = xfrm.find(qn("a:chOff"))
            ch_ext = xfrm.find(qn("a:chExt"))
            off_x, off_y = int(off.get("x")), int(off.get("y"))
            cx, cy = int(ext.get("cx")), int(ext.get("cy"))
            ch_x, ch_y = int(ch_off.get("x")), int(ch_off.get("y"))
            ch_cx, ch_cy = int(ch_ext.get("cx")), int(ch_ext.get("cy"))
        except (AttributeError, TypeError, ValueError):
            return parent
        scale_x = cx / ch_cx if ch_cx else 1.0
        scale_y = cy / ch_cy if ch_cy else 1.0
        return cls(off_x, off_y, ch_x, ch_y, scale_x, scale_y, parent)

    def apply(self, left, top, width, height):
        left = self.off_x + (left - self.ch_off_x) * self.scale_x
        top = self.off_y + (top - self.ch_off_y) * self.scale_y
        width = width * self.scale_x
        height = height * self.scale_y
        if self.parent is not None:
            return self.parent.apply(left, top, width, height)
        return left, top, width, height


# --------------------------------------------------------------------------
# helpers
# --------------------------------------------------------------------------


def _slide_background(slide, theme: Theme) -> str:
    """Slide background, falling back through layout and master."""
    for source in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            fill = source.background.fill
            if fill.type is None:
                continue
            if getattr(fill.type, "name", "") != "SOLID":
                continue
            color = theme.resolve(fill.fore_color)
            if color:
                return color
        except (AttributeError, ValueError, NotImplementedError):
            continue
    return "#FFFFFF"


def _inherited_geometry(shape):
    """A placeholder with no explicit position inherits it from the layout."""
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        idx = shape.placeholder_format.idx
        layout = shape.part.slide.slide_layout
    except (AttributeError, ValueError):
        return None
    for candidate in getattr(layout, "placeholders", []):
        try:
            if candidate.placeholder_format.idx != idx:
                continue
        except (AttributeError, ValueError):
            continue
        if candidate.left is None or not candidate.width:
            return None
        return candidate.left, candidate.top, candidate.width, candidate.height
    return None


def _default_size(shape) -> float:
    if getattr(shape, "is_placeholder", False):
        try:
            return DEFAULT_SIZES.get(shape.placeholder_format.type, DEFAULT_SIZE)
        except (AttributeError, ValueError):
            return DEFAULT_SIZE
    return DEFAULT_SIZE


def _bullet_of(paragraph, *, inherited: LevelDefaults) -> str:
    """Bullet for a paragraph: explicit first, then whatever it inherits.

    Guessing from indent level used to invent bullets on plain text boxes.
    Now the layout/master answers the question, and "no bullet" is a real
    answer rather than the absence of one.
    """
    p_pr = paragraph._p.find(qn("a:pPr"))
    if p_pr is not None:
        if p_pr.find(qn("a:buNone")) is not None:
            return ""
        bu_char = p_pr.find(qn("a:buChar"))
        if bu_char is not None:
            return bu_char.get("char") or "•"
        if p_pr.find(qn("a:buAutoNum")) is not None:
            return "1."
    return inherited.bullet or ""


# OOXML's default rounded-rectangle adjustment: radius = 1/6 of the short side.
DEFAULT_ROUND_ADJ = 16667


def _corner_radius(shape, short_side: float) -> float:
    """Corner radius in pixels, from the shape's own ``adj`` guide.

    A hardcoded radius makes every rounded box look the same; PowerPoint stores
    the real value as a fraction of the shape's shorter side.
    """
    adj = DEFAULT_ROUND_ADJ
    try:
        av_lst = shape._element.find(qn("p:spPr")).find(qn("a:prstGeom")).find(qn("a:avLst"))
        for guide in av_lst.findall(qn("a:gd")):
            if guide.get("name") == "adj":
                adj = int(guide.get("fmla", "val 16667").split()[-1])
                break
    except (AttributeError, TypeError, ValueError, IndexError):
        pass
    return max(0.0, short_side * adj / 100000)


# XL_CHART_TYPE member name -> our render kind. Unlisted types fall back to the
# kind read off the plot group, and only then to a labelled placeholder — being
# drawn as something they are not is the one outcome to avoid.
CHART_KIND_MAP = {
    "COLUMN_CLUSTERED": ChartKind.COLUMN,
    "COLUMN_STACKED": ChartKind.COLUMN_STACKED,
    "COLUMN_STACKED_100": ChartKind.COLUMN_STACKED,
    "BAR_CLUSTERED": ChartKind.BAR,
    "BAR_STACKED": ChartKind.BAR_STACKED,
    "BAR_STACKED_100": ChartKind.BAR_STACKED,
    "LINE": ChartKind.LINE,
    "LINE_MARKERS": ChartKind.LINE_MARKERS,
    "LINE_STACKED": ChartKind.LINE,
    "LINE_STACKED_100": ChartKind.LINE,
    "LINE_MARKERS_STACKED": ChartKind.LINE_MARKERS,
    "LINE_MARKERS_STACKED_100": ChartKind.LINE_MARKERS,
    "PIE": ChartKind.PIE,
    "PIE_EXPLODED": ChartKind.PIE,
    "PIE_OF_PIE": ChartKind.PIE,
    "BAR_OF_PIE": ChartKind.PIE,
    "DOUGHNUT": ChartKind.DOUGHNUT,
    "DOUGHNUT_EXPLODED": ChartKind.DOUGHNUT,
    "AREA": ChartKind.AREA,
    "AREA_STACKED": ChartKind.AREA_STACKED,
    "AREA_STACKED_100": ChartKind.AREA_STACKED,
    "XY_SCATTER": ChartKind.SCATTER,
    "XY_SCATTER_LINES": ChartKind.SCATTER,
    "XY_SCATTER_LINES_NO_MARKERS": ChartKind.SCATTER,
    "XY_SCATTER_SMOOTH": ChartKind.SCATTER,
    "XY_SCATTER_SMOOTH_NO_MARKERS": ChartKind.SCATTER,
    "BUBBLE": ChartKind.SCATTER,
    "BUBBLE_THREE_D_EFFECT": ChartKind.SCATTER,
    # 3D types plot the same data; the depth is a look, not a third dimension.
    "THREE_D_COLUMN": ChartKind.COLUMN,
    "THREE_D_COLUMN_CLUSTERED": ChartKind.COLUMN,
    "THREE_D_COLUMN_STACKED": ChartKind.COLUMN_STACKED,
    "THREE_D_COLUMN_STACKED_100": ChartKind.COLUMN_STACKED,
    "THREE_D_BAR_CLUSTERED": ChartKind.BAR,
    "THREE_D_BAR_STACKED": ChartKind.BAR_STACKED,
    "THREE_D_BAR_STACKED_100": ChartKind.BAR_STACKED,
    "THREE_D_LINE": ChartKind.LINE,
    "THREE_D_AREA": ChartKind.AREA,
    "THREE_D_AREA_STACKED": ChartKind.AREA_STACKED,
    "THREE_D_AREA_STACKED_100": ChartKind.AREA_STACKED,
    "THREE_D_PIE": ChartKind.PIE,
    "THREE_D_PIE_EXPLODED": ChartKind.PIE,
    # Cone / cylinder / pyramid are column and bar charts wearing a shape.
    **{
        f"{prefix}_{suffix}": kind
        for prefix in ("CONE", "CYLINDER", "PYRAMID")
        for suffix, kind in (
            ("COL", ChartKind.COLUMN),
            ("COL_CLUSTERED", ChartKind.COLUMN),
            ("COL_STACKED", ChartKind.COLUMN_STACKED),
            ("COL_STACKED_100", ChartKind.COLUMN_STACKED),
            ("BAR_CLUSTERED", ChartKind.BAR),
            ("BAR_STACKED", ChartKind.BAR_STACKED),
            ("BAR_STACKED_100", ChartKind.BAR_STACKED),
        )
    },
}

# Names whose type is already 3D, for decks whose plot group we could not read.
THREE_D_PREFIXES = ("THREE_D_", "CONE_", "CYLINDER_", "PYRAMID_")


def _is_three_d(type_name: str) -> bool:
    return type_name.startswith(THREE_D_PREFIXES) or type_name == "BUBBLE_THREE_D_EFFECT"


LEGEND_POSITIONS = {"BOTTOM": "bottom", "TOP": "top", "RIGHT": "right", "LEFT": "left"}


def _chart_type_name(chart) -> str:
    try:
        return getattr(chart.chart_type, "name", "") or ""
    except (AttributeError, ValueError):
        return ""


def _chart_categories(chart) -> list[str]:
    try:
        return [str(c) for c in chart.plots[0].categories]
    except (AttributeError, IndexError, ValueError):
        return []


def _chart_title(chart) -> str:
    try:
        if not chart.has_title:
            return ""
        return chart.chart_title.text_frame.text.strip()
    except (AttributeError, ValueError):
        return ""


def _legend_position(chart) -> str:
    try:
        if not chart.has_legend:
            return "bottom"
        name = getattr(chart.legend.position, "name", "") or ""
    except (AttributeError, ValueError):
        return "bottom"
    return LEGEND_POSITIONS.get(name, "bottom")


def _has_data_labels(chart) -> bool:
    try:
        return bool(chart.plots[0].has_data_labels)
    except (AttributeError, IndexError, ValueError, NotImplementedError):
        return False


def _has_gridlines(chart) -> bool:
    try:
        return bool(chart.value_axis.has_major_gridlines)
    except (AttributeError, ValueError, NotImplementedError):
        # Pie charts have no value axis; gridlines are meaningless there anyway.
        return False


def _axis_scale(chart, attribute: str) -> float | None:
    try:
        value = getattr(chart.value_axis, attribute)
    except (AttributeError, ValueError, NotImplementedError):
        return None
    return None if value is None else float(value)


def _style_value(style: TableStyle | None, field: str, fallback: str) -> str:
    """The deck's own value when it defines one, else the approximation."""
    value = getattr(style, field, None) if style is not None else None
    return value or fallback


def _table_style_id(table) -> str | None:
    tbl_pr = table._tbl.find(qn("a:tblPr"))
    if tbl_pr is None:
        return None
    node = tbl_pr.find(qn("a:tableStyleId"))
    return node.text if node is not None and node.text else None


def _has_header(table) -> bool:
    """Whether the deck asked for a banded header row (``firstRow`` in tblPr)."""
    try:
        tbl_pr = table._tbl.find(qn("a:tblPr"))
    except AttributeError:
        return True
    if tbl_pr is None:
        return True
    return tbl_pr.get("firstRow", "1") in ("1", "true")


def _first_set(*values):
    """First value that was actually specified; the last one is the fallback."""
    for value in values:
        if value is not None:
            return value
    return None


LEADING_BULLETS = "•·‣▪◦-*–—"


def _strip_leading_bullet(text: str) -> str:
    stripped = text.lstrip()
    if stripped[:1] in LEADING_BULLETS:
        return stripped[1:].lstrip()
    return text


def _line_spacing(paragraph, runs: list[Run]) -> float:
    """PowerPoint's two line-spacing kinds, both as a CSS line-height multiple.

    ``<a:spcPct>`` gives a multiple and python-pptx returns a float; ``<a:spcPts>``
    gives an absolute height and returns a ``Length``. **Length subclasses int**,
    so it has to be tested first — an ``isinstance(x, int)`` check passes for
    both, and returning a Length unchanged hands the renderer raw EMU as a
    multiple. 21pt of spacing then becomes a line-height of 266700, which pushes
    the text millions of pixels below its own box: it renders, and every
    paragraph using absolute spacing silently disappears.

    An absolute height is divided by the paragraph's own font size, since that
    is what a multiple is relative to.
    """
    spacing = paragraph.line_spacing
    if spacing is None:
        return DEFAULT_LINE_SPACING
    if isinstance(spacing, Length):
        size_pt = next((r.size_pt for r in runs if r.size_pt), None) or DEFAULT_SIZE
        # Bounded because a corrupt value should cost this paragraph's spacing,
        # not the whole slide's layout.
        return max(0.5, min(5.0, spacing.pt / size_pt))
    return float(spacing)


def _length_to_pt(value) -> float:
    if value is None:
        return 0.0
    try:
        return float(value.pt)
    except AttributeError:
        return float(value) / EMU_PER_POINT


def _emu_to_pt(value, default: float) -> float:
    if value is None:
        return default
    try:
        return float(value.pt)
    except AttributeError:
        return float(value) / EMU_PER_POINT
