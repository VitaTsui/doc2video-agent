"""Placeholder text inheritance — the `<a:lstStyle>` chain.

Real decks put almost nothing on the slide itself: a body placeholder's font
size, colour, alignment and bullet character all come from the layout, the
master, or the master's `titleStyle` / `bodyStyle` / `otherStyle` blocks. Reading
only the slide gives you the *text* but not its *appearance*, which is why the
first version of the renderer had to guess sizes from placeholder type.

Resolution order, least specific first (later wins):

1. presentation ``<p:defaultTextStyle>``
2. master ``<p:txStyles>`` — title / body / other, chosen by placeholder type
3. the matching placeholder on the master
4. the matching placeholder on the layout
5. the shape's own ``<a:lstStyle>``
6. properties set directly on the paragraph / run (handled by the extractor)
"""

from __future__ import annotations

from dataclasses import dataclass, replace

from pptx.oxml.ns import qn

from ...core.logging import get_logger
from .theme import Theme

log = get_logger(__name__)

MAX_LEVELS = 9

# OOXML alignment codes -> the values our model uses.
ALIGN_CODES = {
    "l": "left",
    "ctr": "center",
    "r": "right",
    "just": "justify",
    "justLow": "justify",
    "dist": "justify",
}

# Placeholder type -> which master text style governs it.
TITLE_PLACEHOLDERS = {"title", "ctrTitle"}
BODY_PLACEHOLDERS = {"body", "subTitle", "obj", "tbl", "chart", "dgm", "media", "clipArt"}


@dataclass(frozen=True)
class LevelDefaults:
    """Inherited run/paragraph properties for one outline level.

    ``None`` means "not specified at this level" so that merging keeps the more
    specific value; the extractor treats a remaining ``None`` as "use the
    hardcoded fallback".
    """

    size_pt: float | None = None
    bold: bool | None = None
    italic: bool | None = None
    underline: bool | None = None
    color: str | None = None
    font: str | None = None
    align: str | None = None
    # "" means an explicit "no bullet"; None means nothing was said.
    bullet: str | None = None

    def merged_with(self, more_specific: LevelDefaults) -> LevelDefaults:
        changes = {
            field: value
            for field, value in vars(more_specific).items()
            if value is not None
        }
        return replace(self, **changes)


class StyleResolver:
    """Resolves inherited text defaults for the shapes of one presentation."""

    def __init__(self, presentation, theme: Theme) -> None:
        self._theme = theme
        self._presentation = presentation
        self._default_style = self._read_default_text_style(presentation)
        self._master_cache: dict[int, dict[str, dict[int, LevelDefaults]]] = {}
        self._shape_cache: dict[int, dict[int, LevelDefaults]] = {}

    # -- public ------------------------------------------------------------
    def defaults_for(self, shape, level: int) -> LevelDefaults:
        """Inherited defaults for ``shape`` at the given outline level."""
        try:
            levels = self._levels_for_shape(shape)
        except Exception as exc:  # malformed package; fall back to nothing
            log.debug("解析继承样式失败：%s", exc)
            return LevelDefaults()
        return levels.get(min(level, MAX_LEVELS - 1), LevelDefaults())

    def font_scale(self, shape) -> float:
        """PowerPoint's autofit shrink factor, as a multiplier on every size."""
        try:
            body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        except (AttributeError, ValueError):
            return 1.0
        if body_pr is None:
            return 1.0
        norm = body_pr.find(qn("a:normAutofit"))
        if norm is None:
            return 1.0
        try:
            return int(norm.get("fontScale", "100000")) / 100000
        except ValueError:
            return 1.0

    # -- chain -------------------------------------------------------------
    def _levels_for_shape(self, shape) -> dict[int, LevelDefaults]:
        key = id(shape._element)
        cached = self._shape_cache.get(key)
        if cached is not None:
            return cached

        chain: list[dict[int, LevelDefaults]] = [self._default_style]

        placeholder = _placeholder_key(shape)
        layout, master = _layout_and_master(shape)

        if master is not None:
            master_styles = self._master_styles(master)
            style_name = _master_style_name(placeholder)
            if style_name and style_name in master_styles:
                chain.append(master_styles[style_name])
            if placeholder is not None:
                match = _matching_placeholder(master, placeholder)
                if match is not None:
                    chain.append(self._read_lst_style(match))

        if layout is not None and placeholder is not None:
            match = _matching_placeholder(layout, placeholder)
            if match is not None:
                chain.append(self._read_lst_style(match))

        chain.append(self._read_lst_style(shape))

        merged: dict[int, LevelDefaults] = {}
        for level in range(MAX_LEVELS):
            value = LevelDefaults()
            for source in chain:
                if level in source:
                    value = value.merged_with(source[level])
            merged[level] = value

        self._shape_cache[key] = merged
        return merged

    def _master_styles(self, master) -> dict[str, dict[int, LevelDefaults]]:
        key = id(master)
        cached = self._master_cache.get(key)
        if cached is not None:
            return cached

        styles: dict[str, dict[int, LevelDefaults]] = {}
        tx_styles = master.element.find(qn("p:txStyles"))
        if tx_styles is not None:
            for name, tag in (
                ("title", "p:titleStyle"),
                ("body", "p:bodyStyle"),
                ("other", "p:otherStyle"),
            ):
                node = tx_styles.find(qn(tag))
                if node is not None:
                    styles[name] = self._parse_levels(node)

        self._master_cache[key] = styles
        return styles

    def _read_default_text_style(self, presentation) -> dict[int, LevelDefaults]:
        try:
            node = presentation.part._element.find(qn("p:defaultTextStyle"))
        except (AttributeError, ValueError):
            return {}
        return self._parse_levels(node) if node is not None else {}

    def _read_lst_style(self, shape) -> dict[int, LevelDefaults]:
        try:
            body = shape.text_frame._txBody
        except (AttributeError, ValueError):
            return {}
        node = body.find(qn("a:lstStyle"))
        return self._parse_levels(node) if node is not None else {}

    # -- parsing -----------------------------------------------------------
    def _parse_levels(self, container) -> dict[int, LevelDefaults]:
        levels: dict[int, LevelDefaults] = {}
        for level in range(MAX_LEVELS):
            node = container.find(qn(f"a:lvl{level + 1}pPr"))
            if node is None:
                continue
            levels[level] = self._parse_level(node)
        return levels

    def _parse_level(self, node) -> LevelDefaults:
        align = ALIGN_CODES.get(node.get("algn") or "")
        bullet = _bullet_of(node)

        def_rpr = node.find(qn("a:defRPr"))
        if def_rpr is None:
            return LevelDefaults(align=align, bullet=bullet)

        size = def_rpr.get("sz")
        return LevelDefaults(
            # OOXML stores point sizes in hundredths.
            size_pt=int(size) / 100 if size and size.isdigit() else None,
            bold=_flag(def_rpr.get("b")),
            italic=_flag(def_rpr.get("i")),
            underline=None if def_rpr.get("u") is None else def_rpr.get("u") != "none",
            color=self._theme.color_element(def_rpr.find(qn("a:solidFill"))),
            font=self._font_of(def_rpr),
            align=align,
            bullet=bullet,
        )

    def _font_of(self, def_rpr) -> str | None:
        # East Asian first: a CJK deck's latin entry is often a theme reference
        # that would render Chinese in a fallback face.
        for tag in ("a:ea", "a:latin"):
            node = def_rpr.find(qn(tag))
            if node is None:
                continue
            resolved = self._theme.font_ref(node.get("typeface"))
            if resolved:
                return resolved
        return None


# --------------------------------------------------------------------------
# helpers
# --------------------------------------------------------------------------


def _flag(value: str | None) -> bool | None:
    if value is None:
        return None
    return value in ("1", "true")


def _bullet_of(node) -> str | None:
    if node.find(qn("a:buNone")) is not None:
        return ""
    bu_char = node.find(qn("a:buChar"))
    if bu_char is not None:
        return bu_char.get("char") or "•"
    if node.find(qn("a:buAutoNum")) is not None:
        return "1."
    return None


def _placeholder_key(shape) -> tuple[str, int | None] | None:
    """The (type, idx) pair used to match a placeholder across layout/master."""
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        fmt = shape.placeholder_format
        ph_type = getattr(fmt.type, "name", "") or ""
        idx = fmt.idx
    except (AttributeError, ValueError):
        return None
    return (_ooxml_ph_type(ph_type), idx)


def _ooxml_ph_type(name: str) -> str:
    """PP_PLACEHOLDER member name -> the OOXML `type` attribute spelling."""
    return {
        "TITLE": "title",
        "CENTER_TITLE": "ctrTitle",
        "SUBTITLE": "subTitle",
        "BODY": "body",
        "OBJECT": "obj",
        "TABLE": "tbl",
        "CHART": "chart",
        "PICTURE": "pic",
        "SLIDE_NUMBER": "sldNum",
        "FOOTER": "ftr",
        "HEADER": "hdr",
        "DATE": "dt",
    }.get(name, name.lower())


def _master_style_name(placeholder: tuple[str, int | None] | None) -> str:
    if placeholder is None:
        return "other"
    ph_type = placeholder[0]
    if ph_type in TITLE_PLACEHOLDERS:
        return "title"
    if ph_type in BODY_PLACEHOLDERS:
        return "body"
    return "other"


def _layout_and_master(shape):
    """The layout and master that this shape's slide inherits from."""
    try:
        slide = shape.part.slide
        layout = slide.slide_layout
        return layout, layout.slide_master
    except (AttributeError, ValueError):
        return None, None


def _matching_placeholder(source, placeholder: tuple[str, int | None]):
    """Find the placeholder on a layout/master that governs this one.

    PowerPoint matches on ``idx`` first and falls back to the type, which is why
    a title placeholder with no idx still inherits from the layout's title.
    """
    ph_type, idx = placeholder
    candidates = list(getattr(source, "placeholders", []))

    if idx is not None:
        for candidate in candidates:
            try:
                if candidate.placeholder_format.idx == idx:
                    return candidate
            except (AttributeError, ValueError):
                continue

    for candidate in candidates:
        try:
            candidate_type = _ooxml_ph_type(
                getattr(candidate.placeholder_format.type, "name", "") or ""
            )
        except (AttributeError, ValueError):
            continue
        if candidate_type == ph_type:
            return candidate
        # A layout's body placeholder governs obj/tbl/chart content too.
        if ph_type in BODY_PLACEHOLDERS and candidate_type in BODY_PLACEHOLDERS:
            return candidate
    return None
